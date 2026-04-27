#!/usr/bin/env python3
"""
기획서 파서: PPT/PPTX/PDF/DOCX에서 텍스트, 표, 이미지를 분리 추출합니다.
에이전트 친화적(토큰 최소화) 형식으로 출력합니다.

이미지 티어:
  Tier 2 - 텍스트 추출 성공 → 코드블록 인라인, 이미지 파일 미저장
  Tier 3 - 시각 해석 필요 → images/ 저장, [IMG: path] 태그로 참조
"""

import argparse
import glob
import hashlib
import json
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path


_BULLETS = {"•", "·", "▪", "▸", "►", "–", "○", "●"}
_SUPPORTED = {".pdf", ".pptx", ".ppt", ".docx"}


def _progress(current: int, total: int, label: str = "", quiet: bool = False):
    if quiet or not sys.stderr.isatty():
        return
    sys.stderr.write(f"\r  {label} {current}/{total}")
    sys.stderr.flush()
    if current == total:
        sys.stderr.write("\n")


def _join_bullets(texts: list) -> list:
    lines = []
    for text in texts:
        for line in text.split("\n"):
            line = line.strip()
            if line:
                lines.append(line)

    result = []
    i = 0
    while i < len(lines):
        if lines[i] in _BULLETS and i + 1 < len(lines):
            result.append(f"- {lines[i + 1]}")
            i += 2
        else:
            result.append(lines[i])
            i += 1
    return result


def _is_meaningful_text(text: str) -> bool:
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    if not lines:
        return False
    total_words = sum(len(l.split()) for l in lines)
    if total_words < 6:
        return False
    avg_line_len = sum(len(l) for l in lines) / len(lines)
    if avg_line_len < 10:
        return False
    body = text.replace(" ", "").replace("\n", "").replace("\r", "")
    if body:
        _ALLOWED = set('.,;:()[]{}\'\"!?-_=+*/\\@#%^&<>·•▪')
        clean = sum(1 for c in body if c.isalnum() or c in _ALLOWED)
        if clean / len(body) < 0.65:
            return False
    return True


def _is_duplicate_of_page(text: str, page_words: set, threshold: float = 0.6) -> bool:
    words = set(text.split())
    if len(words) < 5:
        return False
    return len(words & page_words) / len(words) >= threshold


def _ocr_image(img_path: Path, warn: bool = False) -> str:
    try:
        import pytesseract
        from PIL import Image, ImageOps, ImageStat
    except ImportError:
        if warn:
            print("  [경고] OCR 불가: pip install pytesseract Pillow && brew install tesseract tesseract-lang")
        return ""
    try:
        img = Image.open(img_path).convert("RGB")
        gray = img.convert("L")
        if ImageStat.Stat(gray).mean[0] < 128:
            img = ImageOps.invert(img)
        w, h = img.size
        img = img.resize((w * 2, h * 2), Image.LANCZOS)
        return pytesseract.image_to_string(img, lang="kor+eng").strip()
    except Exception:
        return ""


def _ocr_from_pixmap(pix, warn: bool = False) -> str:
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        tmp = Path(f.name)
    try:
        pix.save(str(tmp))
        return _ocr_image(tmp, warn=warn)
    finally:
        tmp.unlink(missing_ok=True)


def _rows_to_markdown(rows: list) -> str:
    if not rows or not rows[0]:
        return ""
    rows = [r for r in rows if any(str(cell).strip() for cell in r)]
    if not rows:
        return ""
    max_cols = max(len(r) for r in rows)
    rows = [list(r) + [""] * (max_cols - len(r)) for r in rows]
    # 셀 내 줄바꿈·파이프 이스케이프
    for i, row in enumerate(rows):
        rows[i] = [str(c).replace("|", "\\|").replace("\n", " ") for c in row]

    header = "| " + " | ".join(rows[0]) + " |"
    separator = "| " + " | ".join("---" for _ in rows[0]) + " |"
    body_lines = ["| " + " | ".join(row) + " |" for row in rows[1:]]
    if body_lines:
        return f"{header}\n{separator}\n" + "\n".join(body_lines)
    return f"{header}\n{separator}"


def _detect_pdf_headings(page, blocks_text: list) -> list:
    try:
        d = page.get_text("dict")
    except Exception:
        return blocks_text

    sizes = []
    heading_candidates = {}

    for block in d.get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            spans = line.get("spans", [])
            text = "".join(s.get("text", "") for s in spans).strip()
            if not text:
                continue
            max_size = max((s.get("size", 0) for s in spans), default=0)
            is_bold = any(s.get("flags", 0) & 16 for s in spans)
            sizes.append(max_size)
            heading_candidates[text] = (max_size, is_bold)

    if len(sizes) < 3:
        return blocks_text

    sizes.sort()
    median = sizes[len(sizes) // 2]

    heading_map = {}
    for text, (size, bold) in heading_candidates.items():
        if size > median * 1.4:
            heading_map[text] = "##"
        elif size > median * 1.15 and bold:
            heading_map[text] = "###"

    result = []
    for line in blocks_text:
        clean = line.lstrip("- ").strip()
        if clean in heading_map:
            result.append(f"{heading_map[clean]} {clean}")
        else:
            result.append(line)
    return result


def _md_header(path: Path, file_type: str, n_sections: int, tier3_images: list, n_tables: int = 0) -> str:
    section_label = "슬라이드" if file_type in ("pptx", "ppt") else ("섹션" if file_type == "docx" else "페이지")
    lines = [
        f"# 문서: {path.name}",
        f"- 형식: {file_type.upper()} | {section_label}: {n_sections}",
    ]
    if n_tables:
        lines.append(f"- 표: {n_tables}개")
    lines.append(f"- Tier3 이미지 (에이전트 확인 필요): {len(tier3_images)}개")
    if tier3_images:
        for p in tier3_images:
            lines.append(f"  - {p}")
    return "\n".join(lines)


# ─── PPTX ───────────────────────────────────────────────────────────────────

def extract_pptx(path: Path, out_dir: Path, use_ocr: bool = False, quiet: bool = False) -> dict:
    try:
        from pptx import Presentation
    except ImportError:
        sys.exit("python-pptx 미설치: pip install python-pptx")

    img_dir = out_dir / "images"
    img_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(path)
    slides_md = []
    tier3_images = []
    table_count = 0
    manifest = {"source": str(path), "type": "pptx", "slides": []}
    total_slides = len(prs.slides)
    errors = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        _progress(slide_idx, total_slides, "슬라이드", quiet)
        try:
            slide_info = {"slide": slide_idx, "title": None, "images": [], "tables": 0}

            title = ""
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
                slide_info["title"] = title

            lines = [f"## [S{slide_idx}]{' ' + title if title else ''}"]

            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        prefix = "  " * para.level + "-" if para.level > 0 else "-"
                        lines.append(f"{prefix} {text}")

            # 표 추출
            for shape in slide.shapes:
                if shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        rows.append([cell.text.strip() for cell in row.cells])
                    md_table = _rows_to_markdown(rows)
                    if md_table:
                        lines.append(f"\n{md_table}")
                        table_count += 1
                        slide_info["tables"] += 1

            page_words = set(" ".join(lines).split())

            img_count = 0
            seen_parts = set()
            for rel in slide.part.rels.values():
                if "image" not in rel.reltype.lower():
                    continue
                try:
                    img_part = rel.target_part
                except Exception:
                    continue
                part_id = id(img_part)
                if part_id in seen_parts:
                    continue
                seen_parts.add(part_id)

                ext = img_part.partname.suffix.lstrip(".")
                if ext.lower() in ("emf", "wmf"):
                    continue

                img_name = f"slide_{slide_idx}_img_{img_count}.{ext}"
                img_path = img_dir / img_name
                img_path.write_bytes(img_part.blob)

                tier = 3
                extracted_text = ""

                if use_ocr:
                    ocr = _ocr_image(img_path, warn=True)
                    if ocr and not _is_duplicate_of_page(ocr, page_words):
                        extracted_text = ocr
                        tier = 2

                if tier == 2:
                    lines.append(f"\n```\n{extracted_text}\n```")
                    img_path.unlink()
                    slide_info["images"].append({"id": img_name, "tier": 2})
                else:
                    lines.append(f"\n[IMG: images/{img_name}]")
                    tier3_images.append(f"images/{img_name}")
                    slide_info["images"].append({"id": img_name, "tier": 3, "path": f"images/{img_name}"})

                img_count += 1

            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"\n> 노트: {notes}")

            slides_md.append("\n".join(lines))
            manifest["slides"].append(slide_info)

        except Exception as e:
            errors.append(f"슬라이드 {slide_idx}: {e}")
            slides_md.append(f"## [S{slide_idx}]\n\n> 추출 실패: {e}")
            manifest["slides"].append({"slide": slide_idx, "error": str(e)})

    header = _md_header(path, "pptx", total_slides, tier3_images, table_count)
    manifest["summary"] = {
        "total_slides": total_slides,
        "tables": table_count,
        "tier3_count": len(tier3_images),
        "tier3_images": tier3_images,
        "errors": errors,
    }
    full_md = header + "\n\n---\n\n" + "\n\n---\n\n".join(slides_md)
    return {"markdown": full_md, "manifest": manifest}


# ─── PDF ────────────────────────────────────────────────────────────────────

def extract_pdf(
    path: Path,
    out_dir: Path,
    min_area_ratio: float = 0.02,
    render_scale: float = 2.0,
    use_ocr: bool = False,
    quiet: bool = False,
) -> dict:
    try:
        import fitz
    except ImportError:
        sys.exit("pymupdf 미설치: pip install pymupdf")

    img_dir = out_dir / "images"
    img_dir.mkdir(parents=True, exist_ok=True)

    doc = fitz.open(str(path))
    pages_md = []
    tier3_images = []
    table_count = 0
    manifest = {"source": str(path), "type": "pdf", "pages": []}
    seen_hashes: set = set()
    total_pages = len(doc)
    errors = []

    for page_idx, page in enumerate(doc, 1):
        _progress(page_idx, total_pages, "페이지", quiet)
        try:
            page_info = {"page": page_idx, "images": [], "tables": 0}

            page_rect = page.rect
            page_area = page_rect.width * page_rect.height

            # 표 영역 감지 (텍스트 추출 전에 먼저 수행)
            table_rects = []
            table_entries = []
            try:
                for table in page.find_tables():
                    rect = fitz.Rect(table.bbox)
                    table_rects.append(rect)
                    cells = table.extract()
                    rows = [[c or "" for c in row] for row in cells]
                    if any(any(str(c).strip() for c in row) for row in rows):
                        table_entries.append((rect.y0, rows))
            except Exception:
                pass

            # 텍스트 블록 추출 — 표 영역과 겹치는 블록 제외
            text_entries = []
            for b in page.get_text("blocks", sort=True):
                if b[6] != 0 or not b[4].strip():
                    continue
                br = fitz.Rect(b[:4])
                overlaps = False
                for tr in table_rects:
                    overlap = br & tr
                    if not overlap.is_empty:
                        block_area = max(br.width * br.height, 1)
                        if (overlap.width * overlap.height) / block_area > 0.5:
                            overlaps = True
                            break
                if not overlaps:
                    text_entries.append((b[1], b[4].strip()))

            # Y좌표 기준으로 텍스트·표 병합
            lines = [f"## [P{page_idx}]"]
            all_items = [(y, "text", data) for y, data in text_entries] + \
                        [(y, "table", data) for y, data in table_entries]
            all_items.sort(key=lambda x: x[0])

            text_buf = []
            for _, kind, data in all_items:
                if kind == "text":
                    text_buf.append(data)
                else:
                    if text_buf:
                        cleaned = _join_bullets(text_buf)
                        cleaned = _detect_pdf_headings(page, cleaned)
                        lines.extend(cleaned)
                        text_buf = []
                    md_table = _rows_to_markdown(data)
                    if md_table:
                        lines.append(f"\n{md_table}")
                        table_count += 1
                        page_info["tables"] += 1
            if text_buf:
                cleaned = _join_bullets(text_buf)
                cleaned = _detect_pdf_headings(page, cleaned)
                lines.extend(cleaned)

            page_words = set(" ".join(lines).split())

            # 이미지
            try:
                img_infos = page.get_image_info(xrefs=True)
            except TypeError:
                img_infos = page.get_image_info()

            seen_xrefs: set = set()
            mat = fitz.Matrix(render_scale, render_scale)
            img_count = 0

            for info in img_infos:
                xref = info.get("xref", 0)
                if xref in seen_xrefs:
                    continue

                bbox = fitz.Rect(info["bbox"])
                if bbox.width * bbox.height < page_area * min_area_ratio:
                    continue
                bbox = bbox & page_rect
                if bbox.is_empty:
                    continue
                seen_xrefs.add(xref)

                pix = page.get_pixmap(matrix=mat, clip=bbox)
                img_hash = hashlib.sha256(pix.tobytes()).hexdigest()
                if img_hash in seen_hashes:
                    continue
                seen_hashes.add(img_hash)

                clip_text = page.get_text("text", clip=bbox).strip()
                tier = 3
                extracted_text = ""

                if clip_text and _is_meaningful_text(clip_text) and not _is_duplicate_of_page(clip_text, page_words):
                    extracted_text = clip_text
                    tier = 2
                else:
                    ocr = _ocr_from_pixmap(pix, warn=use_ocr)
                    if ocr and _is_meaningful_text(ocr) and not _is_duplicate_of_page(ocr, page_words):
                        extracted_text = ocr
                        tier = 2

                if tier == 2:
                    lines.append(f"\n```\n{extracted_text}\n```")
                    page_info["images"].append({"id": f"page_{page_idx}_img_{img_count}", "tier": 2})
                else:
                    img_name = f"page_{page_idx}_img_{img_count}.png"
                    img_path = img_dir / img_name
                    pix.save(str(img_path))
                    lines.append(f"\n[IMG: images/{img_name}]")
                    tier3_images.append(f"images/{img_name}")
                    page_info["images"].append({"id": img_name, "tier": 3, "path": f"images/{img_name}"})

                img_count += 1

            pages_md.append("\n".join(lines))
            manifest["pages"].append(page_info)

        except Exception as e:
            errors.append(f"페이지 {page_idx}: {e}")
            pages_md.append(f"## [P{page_idx}]\n\n> 추출 실패: {e}")
            manifest["pages"].append({"page": page_idx, "error": str(e)})

    doc.close()

    total_imgs = sum(len(p.get("images", [])) for p in manifest["pages"])
    header = _md_header(path, "pdf", total_pages, tier3_images, table_count)
    manifest["summary"] = {
        "total_pages": total_pages,
        "total_images": total_imgs,
        "tables": table_count,
        "tier3_count": len(tier3_images),
        "tier3_images": tier3_images,
        "errors": errors,
    }
    full_md = header + "\n\n---\n\n" + "\n\n---\n\n".join(pages_md)
    return {"markdown": full_md, "manifest": manifest}


# ─── DOCX ──────────────────────────────────────────────────────────────────

def extract_docx(path: Path, out_dir: Path, use_ocr: bool = False, quiet: bool = False) -> dict:
    try:
        from docx import Document
    except ImportError:
        sys.exit("python-docx 미설치: pip install python-docx")

    from docx.oxml.ns import qn

    img_dir = out_dir / "images"
    img_dir.mkdir(parents=True, exist_ok=True)

    doc = Document(path)
    lines = []
    tier3_images = []
    table_count = 0
    img_count = 0
    section_count = 0
    manifest = {"source": str(path), "type": "docx", "sections": []}
    errors = []

    # XML 순서대로 단락·표 순회
    body = doc.element.body
    para_map = {id(p._element): p for p in doc.paragraphs}
    table_map = {id(t._element): t for t in doc.tables}
    children = list(body)
    total = len(children)

    for idx, child in enumerate(children):
        if idx % 20 == 0:
            _progress(min(idx + 1, total), total, "요소", quiet)
        try:
            if child.tag == qn("w:p") and id(child) in para_map:
                para = para_map[id(child)]
                text = para.text.strip()
                if not text:
                    continue
                style_name = para.style.name if para.style else ""

                if "Heading 1" in style_name or "제목 1" in style_name:
                    section_count += 1
                    lines.append(f"\n## {text}")
                    manifest["sections"].append({"section": section_count, "title": text})
                elif "Heading 2" in style_name or "제목 2" in style_name:
                    lines.append(f"\n### {text}")
                elif "Heading" in style_name or "제목" in style_name:
                    lines.append(f"\n#### {text}")
                elif "List" in style_name or "목록" in style_name:
                    lines.append(f"- {text}")
                else:
                    lines.append(text)

            elif child.tag == qn("w:tbl") and id(child) in table_map:
                table = table_map[id(child)]
                rows = []
                for row in table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                md_table = _rows_to_markdown(rows)
                if md_table:
                    lines.append(f"\n{md_table}")
                    table_count += 1

        except Exception as e:
            errors.append(f"요소 {idx}: {e}")

    _progress(total, total, "요소", quiet)

    # 이미지 추출
    page_words = set(" ".join(lines).split())
    for rel in doc.part.rels.values():
        if "image" not in rel.reltype.lower():
            continue
        try:
            img_part = rel.target_part
        except Exception:
            continue

        ext = Path(img_part.partname).suffix.lstrip(".")
        if ext.lower() in ("emf", "wmf"):
            continue

        img_name = f"docx_img_{img_count}.{ext}"
        img_path = img_dir / img_name
        img_path.write_bytes(img_part.blob)

        tier = 3
        if use_ocr:
            ocr = _ocr_image(img_path, warn=True)
            if ocr and not _is_duplicate_of_page(ocr, page_words):
                lines.append(f"\n```\n{ocr}\n```")
                img_path.unlink()
                tier = 2

        if tier == 3:
            lines.append(f"\n[IMG: images/{img_name}]")
            tier3_images.append(f"images/{img_name}")

        img_count += 1

    n_sections = section_count or 1
    header = _md_header(path, "docx", n_sections, tier3_images, table_count)
    manifest["summary"] = {
        "total_sections": n_sections,
        "tables": table_count,
        "total_images": img_count,
        "tier3_count": len(tier3_images),
        "tier3_images": tier3_images,
        "errors": errors,
    }
    full_md = header + "\n\n---\n\n" + "\n".join(lines)
    return {"markdown": full_md, "manifest": manifest}


# ─── PPT 변환 ────────────────────────────────────────────────────────────────

def convert_ppt_to_pptx(ppt_path: Path) -> Path:
    libreoffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not libreoffice:
        sys.exit(".ppt 변환에 LibreOffice 필요\n설치: brew install --cask libreoffice")

    tmp_dir = Path(tempfile.mkdtemp())
    result = subprocess.run(
        [libreoffice, "--headless", "--convert-to", "pptx", "--outdir", str(tmp_dir), str(ppt_path)],
        capture_output=True, text=True,
    )
    if result.returncode != 0:
        sys.exit(f"PPT 변환 실패:\n{result.stderr}")
    converted = list(tmp_dir.glob("*.pptx"))
    if not converted:
        sys.exit("변환된 파일을 찾을 수 없습니다.")
    return converted[0]


def _build_guide(source: Path, section: str, n: int, tier3_images: list) -> str:
    if tier3_images:
        img_list = "\n".join(f"- {p}" for p in tier3_images)
        img_section = f"\n시각 확인 필요 이미지:\n{img_list}"
    else:
        img_section = "\n이미지 없음 — content.md만 읽으면 됩니다."

    return (
        f"원본: {source.name} | {section}: {n}\n"
        f"읽기: content.md → `[IMG: path]` 태그 위치에서 해당 이미지 로드"
        f"{img_section}"
    )


# ─── 진입점 ──────────────────────────────────────────────────────────────────

def run(input_path: Path, out_dir: Path, min_area_ratio: float, use_ocr: bool, quiet: bool = False) -> None:
    suffix = input_path.suffix.lower()

    if suffix == ".pptx":
        result = extract_pptx(input_path, out_dir, use_ocr=use_ocr, quiet=quiet)
        key = "slides"
    elif suffix == ".ppt":
        if not quiet:
            print("PPT → PPTX 변환 중...")
        result = extract_pptx(convert_ppt_to_pptx(input_path), out_dir, use_ocr=use_ocr, quiet=quiet)
        key = "slides"
    elif suffix == ".pdf":
        result = extract_pdf(input_path, out_dir, min_area_ratio=min_area_ratio, use_ocr=use_ocr, quiet=quiet)
        key = "pages"
    elif suffix == ".docx":
        result = extract_docx(input_path, out_dir, use_ocr=use_ocr, quiet=quiet)
        key = "sections"
    else:
        sys.exit(f"지원하지 않는 형식: {suffix}  (지원: {', '.join(_SUPPORTED)})")

    md_path = out_dir / "content.md"
    md_path.write_text(result["markdown"], encoding="utf-8")

    manifest = result["manifest"]
    manifest_path = out_dir / "manifest.json"
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    s = manifest["summary"]
    tier3 = s.get("tier3_images", [])
    total_key = f"total_{key}"
    n = s.get(total_key, s.get("total_pages", s.get("total_sections", 0)))
    section_map = {"slides": "슬라이드", "pages": "페이지", "sections": "섹션"}
    section = section_map.get(key, "섹션")

    guide = _build_guide(input_path, section, n, tier3)
    (out_dir / "GUIDE.md").write_text(guide, encoding="utf-8")

    errs = s.get("errors", [])
    tables = s.get("tables", 0)

    if not quiet:
        print(f"\n추출 완료: {out_dir}")
        print(f"  {section}: {n}개")
        if tables:
            print(f"  표: {tables}개")
        print(f"  Tier3 이미지 (에이전트 확인): {s['tier3_count']}개")
        print(f"  텍스트: {md_path.name}")
        print(f"  매니페스트: {manifest_path.name}")
        if errs:
            print(f"  경고: {len(errs)}건 오류 발생 (manifest.json 참조)")


def _resolve_files(patterns: list) -> list:
    files = []
    for p in patterns:
        expanded = glob.glob(str(p))
        if expanded:
            for f in sorted(expanded):
                fp = Path(f).resolve()
                if fp.is_file() and fp.suffix.lower() in _SUPPORTED:
                    files.append(fp)
        else:
            fp = Path(p).resolve()
            if fp.exists():
                files.append(fp)
    return files


def main():
    parser = argparse.ArgumentParser(
        description="PPT/PPTX/PDF/DOCX → 에이전트 친화적 마크다운 추출",
    )
    parser.add_argument("files", nargs="+", type=str, metavar="FILE",
                        help="입력 파일 (glob 패턴 지원, 예: *.pdf)")
    parser.add_argument("-o", "--output", type=Path, default=None,
                        help="출력 디렉터리 (단일 파일) 또는 상위 디렉터리 (다중 파일)")
    parser.add_argument("--min-img-ratio", type=float, default=0.02, metavar="RATIO")
    parser.add_argument("--ocr", action="store_true", help="OCR 경고 메시지 활성화")
    parser.add_argument("-q", "--quiet", action="store_true", help="진행률·결과 출력 억제")
    args = parser.parse_args()

    input_files = _resolve_files(args.files)
    if not input_files:
        sys.exit(f"파일 없음: {args.files}")

    if len(input_files) == 1:
        f = input_files[0]
        out_dir = args.output or f.parent / f"{f.stem}_output"
        out_dir.mkdir(parents=True, exist_ok=True)
        run(f, out_dir, min_area_ratio=args.min_img_ratio, use_ocr=args.ocr, quiet=args.quiet)
    else:
        base_dir = args.output or Path.cwd() / "parsed_output"
        base_dir.mkdir(parents=True, exist_ok=True)
        if not args.quiet:
            print(f"배치 모드: {len(input_files)}개 파일\n")
        for i, f in enumerate(input_files, 1):
            if not args.quiet:
                print(f"[{i}/{len(input_files)}] {f.name}")
            out_dir = base_dir / f.stem
            out_dir.mkdir(parents=True, exist_ok=True)
            try:
                run(f, out_dir, min_area_ratio=args.min_img_ratio, use_ocr=args.ocr, quiet=args.quiet)
            except SystemExit:
                raise
            except Exception as e:
                print(f"  오류: {e}", file=sys.stderr)
        if not args.quiet:
            print(f"\n전체 완료: {base_dir}")


if __name__ == "__main__":
    main()
