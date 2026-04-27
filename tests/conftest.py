"""프로그래밍 방식으로 테스트 픽스처를 생성 (외부 파일 불필요)"""
import json
import tempfile
from pathlib import Path

import pytest


@pytest.fixture
def tmp_out(tmp_path):
    return tmp_path / "output"


@pytest.fixture
def sample_pdf(tmp_path):
    """최소 PDF 파일 생성 (텍스트 + 테이블 포함)"""
    import fitz

    path = tmp_path / "sample.pdf"
    doc = fitz.open()

    # 페이지 1: 큰 제목 + 본문
    page = doc.new_page()
    page.insert_text((72, 80), "문서 제목", fontsize=24, fontname="helv")
    page.insert_text((72, 130), "본문 텍스트입니다. 이것은 테스트용 문서의 첫 번째 페이지입니다.", fontsize=11, fontname="helv")
    page.insert_text((72, 160), "추가 내용이 여기에 들어갑니다. 충분한 단어 수를 확보합니다.", fontsize=11, fontname="helv")

    # 페이지 2: 본문
    page2 = doc.new_page()
    page2.insert_text((72, 80), "두 번째 페이지", fontsize=16, fontname="helv")
    page2.insert_text((72, 130), "두 번째 페이지의 내용입니다. 다양한 텍스트가 포함되어 있습니다.", fontsize=11, fontname="helv")

    doc.save(str(path))
    doc.close()
    return path


@pytest.fixture
def sample_pptx(tmp_path):
    """최소 PPTX 파일 생성 (텍스트 + 테이블 포함)"""
    from pptx import Presentation
    from pptx.util import Inches

    path = tmp_path / "sample.pptx"
    prs = Presentation()

    # 슬라이드 1: 제목 + 텍스트
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "테스트 슬라이드"
    slide.placeholders[1].text = "본문 텍스트입니다."

    # 슬라이드 2: 테이블
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    rows, cols = 3, 3
    table_shape = slide2.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(6), Inches(2))
    table = table_shape.table
    data = [["항목", "수량", "비고"], ["A", "10", "정상"], ["B", "20", "확인필요"]]
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    prs.save(str(path))
    return path


@pytest.fixture
def sample_docx(tmp_path):
    """최소 DOCX 파일 생성 (헤딩 + 단락 + 테이블)"""
    from docx import Document

    path = tmp_path / "sample.docx"
    doc = Document()
    doc.add_heading("문서 제목", level=1)
    doc.add_paragraph("본문 텍스트입니다. 이것은 테스트용 DOCX 문서입니다.")
    doc.add_heading("하위 제목", level=2)
    doc.add_paragraph("두 번째 섹션의 내용입니다.")

    table = doc.add_table(rows=3, cols=2)
    table.cell(0, 0).text = "이름"
    table.cell(0, 1).text = "값"
    table.cell(1, 0).text = "항목A"
    table.cell(1, 1).text = "100"
    table.cell(2, 0).text = "항목B"
    table.cell(2, 1).text = "200"

    doc.save(str(path))
    return path
